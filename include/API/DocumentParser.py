# coding:utf-8
# 做其他事情前，记得先把下面几个模块装了：python-docx python-pptx pdfplumber，
import os
from docx import Document
import xlrd
from pptx import Presentation
import pdfplumber


# pdf读取示例 https://blog.csdn.net/weixin_42305022/article/details/105776403
#

class WordParser:
    def __init__(self):
        pass

    def parse(self, pth_doc):
        _doc = []
        document = Document(pth_doc)  # 打开文件demo.docx
        for paragraph in document.paragraphs:
            _doc.append(paragraph.text)
        return _doc


class ExcelParser:
    def __init__(self):
        pass

    def parse(self, pth_doc):
        # 文件名以及路径，如果路径或者文件名有中文给前面加一个表示原生字符。
        book = xlrd.open_workbook(pth_doc)
        table = book.sheets()[0]

        nrows = table.nrows
        ncols = table.ncols

        _doc = []
        for row in range(nrows):
            vcols = [str(table.cell_value(row, col)) for col in range(ncols)]
            _doc.append(', '.join(vcols))
        return _doc


class PPTParser:
    def __init__(self):
        pass

    def parse(self, pth_doc):
        prs = Presentation(pth_doc)

        _doc = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        _doc.append(run.text)
        return _doc


class PDFParser:
    def __init__(self):
        pass

    def parse(self, pth_doc):
        pdf_plumber = pdfplumber.open(pth_doc)
        # pdf_camelot = camelot.read_pdf(self.doc_pth, pages='', flavor='stream')
        _pages = pdf_plumber.pages

        _doc = []
        for num_page, _page in enumerate(_pages):
            text_page = ''
            try:
                text_page = _page.extract_text()
                _doc.append(text_page)
            except Exception as e:
                print(e)
        return _doc


class DocumentParser:
    # 双下划綫__foo__: 定义的是特殊方法，一般是系统定义名字 ，类似 __init__() 之类的。
    #
    # _foo: 以单下划线开头的表示的是 protected 类型的变量，即保护类型只能允许其本身与子类进行访问，不能用于 from module import *
    #
    # __foo: 双下划线的表示的是私有类型(private)的变量, 只能是允许这个类本身进行访问了。
    def __init__(self):
        self._parser_word = WordParser()
        self._parser_excel = ExcelParser()
        self._parser_ppt = PPTParser()
        self._parser_pdf = PDFParser()
        self.parsers = {
            'docx': self._parser_word,
            'xlsx': self._parser_excel,
            'pptx': self._parser_ppt,
            'pdf': self._parser_pdf
        }

    def parse(self, pth_doc, _type):
        # 下面这句代表返回path最后的文件名。如果path以／或\结尾，那么就会返回空值。即os.path.split(path)的第二个元素。
        file_name = os.path.basename(pth_doc)
        # 分离文件名与扩展名
        _, ext = os.path.splitext(file_name)
        # 两个下划线开头，声明该属性为私有，不能在类的外部被使用或直接访问。
        __type = ext[1:]

        _doc = self.parsers[_type].parse(pth_doc)
        return _doc

# 测试遗留代码
# test = ExcelParser()
# test = WordParser()
# test = PPTParser()
# test = PDFParser()

# test = DocumentParser()
# op = test.parse('A:/mycat/测试.pptx','pptx')
# 上面这行的地址 1注意文件起名不要t开头，好像会变成转义字符啥的 2 斜杠是现在这样朝左的，不要弄反了
# print(op)
